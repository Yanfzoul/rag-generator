#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Batch-convert documents (PDF + Office/LibreOffice) to plain text using
`unstructured.partition.auto.partition`. Preserves directory structure.

Supported suffixes:
  .pdf, .doc, .docx, .ppt, .pptx, .xls, .xlsx, .odt, .ods, .odp, .rtf

Usage:
  py ingest/convert_docs_to_text.py --src ./data/docs --dst ./data/docs_txt
  # (optional) limit to patterns: --glob "**/*.pdf" "**/*.docx"
  # (optional) overwrite existing outputs: --overwrite

Notes:
  - Skips files that have an up-to-date .txt already in dst unless --overwrite.
  - Writes a simple JSONL report to <dst>/convert_report.jsonl
"""

from __future__ import annotations
import argparse, json, os, sys
from pathlib import Path
import yaml
from typing import Iterable
from types import SimpleNamespace
from collections import defaultdict
import re

DOC_SUFFIXES = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".odt", ".ods", ".odp", ".rtf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}


def iter_files(src: Path, globs: list[str] | None) -> Iterable[Path]:
    if globs:
        seen: set[Path] = set()
        for g in globs:
            for p in src.glob(g):
                if p.is_file() and p.suffix.lower() in DOC_SUFFIXES:
                    seen.add(p)
        yield from sorted(seen)
        return
    for p in src.rglob('*'):
        if p.is_file() and p.suffix.lower() in DOC_SUFFIXES:
            yield p


def convert_one(
    src_root: Path,
    dst_root: Path,
    p: Path,
    overwrite: bool = False,
    enable_ocr: bool = False,
    ocr_lang: str | None = None,
    use_tesseract: bool = True,
    ocr_conf_threshold: float = 0.0,
    ocr_image_conf_threshold: float = 0.6,
) -> dict:
    rel = p.relative_to(src_root)
    out = (dst_root / rel).with_suffix('.txt')
    out.parent.mkdir(parents=True, exist_ok=True)

    # Skip if up-to-date
    try:
        if out.exists() and not overwrite:
            if out.stat().st_mtime >= p.stat().st_mtime:
                return {"path": str(p), "out": str(out), "status": "skipped"}
    except Exception:
        pass

    try:
        suffix = p.suffix.lower()
        if suffix == ".docx":
            elements = convert_docx_with_mammoth(str(p))
        elif suffix == ".pptx":
            elements = convert_pptx_with_pythonpptx(str(p))
        elif suffix in {".xls", ".xlsx"}:
            elements = convert_excel_with_pandas(str(p))
        elif suffix == ".pdf":
            native = extract_pdf_text_with_fitz(str(p))
            if native:
                elements = [SimpleNamespace(text=native)]
            else:
                pdf_conf = ocr_conf_threshold if ocr_conf_threshold > 0 else 0.0
                elements = partition_with_ocr(str(p), enable_ocr, ocr_lang, use_tesseract, pdf_conf, ocr_image_conf_threshold)
        elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}:
            conf = ocr_conf_threshold if ocr_conf_threshold > 0 else ocr_image_conf_threshold
            text = ocr_image(str(p), ocr_lang or "eng", use_tesseract, conf)
            elements = [SimpleNamespace(text=text)]
        else:
            conf = ocr_conf_threshold if ocr_conf_threshold > 0 else 0.0
            elements = partition_with_ocr(str(p), enable_ocr, ocr_lang, use_tesseract, conf, ocr_image_conf_threshold)
        text = "\n\n".join(getattr(el, "text", "") for el in elements if getattr(el, "text", ""))
        text = text.strip()
        if not text:
            out.write_text("", encoding='utf-8')
            return {"path": str(p), "out": str(out), "status": "empty"}
        out.write_text(text, encoding='utf-8', errors='ignore')
        return {"path": str(p), "out": str(out), "status": "ok", "chars": len(text)}
    except Exception as e:
        return {"path": str(p), "error": str(e), "status": "error"}


def partition_with_ocr(path: str, enable_ocr: bool, ocr_lang: str | None, use_tesseract: bool, ocr_conf_threshold: float, ocr_image_conf_threshold: float):
    suffix = Path(path).suffix.lower()
    if enable_ocr and suffix == ".pdf":
        try:
            from unstructured.partition.pdf import partition_pdf

            languages = _parse_lang(ocr_lang)
            return partition_pdf(
                filename=path,
                strategy="ocr_only",
                infer_table_structure=True,
                languages=languages,
            )
        except Exception as err:
            print(f"[convert] OCR partition failed for {path}: {err}; trying PyMuPDF fallback", file=sys.stderr)
            conf = ocr_conf_threshold if ocr_conf_threshold > 0 else ocr_image_conf_threshold
            fallback = ocr_pdf_with_pymupdf(path, ocr_lang or "eng", use_tesseract=use_tesseract, conf_threshold=conf)
            if fallback is not None:
                return fallback
            print("[convert] PyMuPDF fallback unavailable or failed; falling back to auto", file=sys.stderr)
    from unstructured.partition.auto import partition

    kwargs = {"filename": path}
    if enable_ocr and suffix == ".pdf":
        kwargs["infer_table_structure"] = True
        kwargs["process_ocr_pdf"] = True
    return partition(**kwargs)


def convert_docx_with_mammoth(path: str):
    try:
        import mammoth
    except ImportError:
        raise RuntimeError("mammoth not installed; run pip install mammoth")
    with open(path, "rb") as f:
        result = mammoth.convert_to_markdown(
            f,
            convert_image=mammoth.images.img_element(lambda image: {"alt_text": image.alt_text or "[image omitted]"})
        )
    text = result.value or ""
    return [SimpleNamespace(text=text)]


def convert_pptx_with_pythonpptx(path: str):
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed; run pip install python-pptx")
    prs = Presentation(path)
    blocks: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        header = f"## Slide {idx}"
        if title:
            header += f": {title}"
        blocks.append(header)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                tbl = shape.table
                rows = []
                for r in tbl.rows:
                    row = [cell.text.strip() for cell in r.cells]
                    rows.append(row)
                blocks.append(render_table(rows))
    text = "\n\n".join(blocks).strip()
    return [SimpleNamespace(text=text)]


def convert_excel_with_pandas(path: str):
    try:
        import pandas as pd
    except ImportError:
        raise RuntimeError("pandas not installed; run pip install pandas")
    try:
        sheets = pd.read_excel(path, sheet_name=None)
    except ImportError as err:
        raise RuntimeError("Missing Excel engine (install openpyxl for .xlsx or xlrd for .xls)") from err
    blocks: list[str] = []
    for name, df in (sheets or {}).items():
        if df is None or df.empty:
            continue
        blocks.append(f"### Sheet: {name}")
        try:
            table_md = df.to_markdown(index=False)
        except Exception:
            table_md = df.to_csv(index=False)
        blocks.append(table_md)
    text = "\n\n".join(blocks).strip()
    return [SimpleNamespace(text=text)] if text else [SimpleNamespace(text="")]


def extract_pdf_text_with_fitz(path: str, min_chars: int = 200) -> str:
    try:
        import fitz  # PyMuPDF
        from markdownify import markdownify as md
    except ImportError:
        return ""
    text_chunks: list[str] = []
    try:
        with fitz.open(path) as doc:
            for page in doc:
                html = page.get_text("xhtml") or page.get_text("html")
                if not html:
                    continue
                link_html = []
                for link in page.get_links():
                    uri = (link or {}).get("uri")
                    rect = (link or {}).get("from")
                    if not uri or not rect:
                        continue
                    snippet = (page.get_textbox(fitz.Rect(rect)) or "").strip()
                    anchor = snippet or uri
                    link_html.append(f'<p><a href="{uri}">{anchor}</a></p>')
                if link_html:
                    html += "\n".join(link_html)
                text_chunks.append(html)
    except Exception:
        return ""
    if not text_chunks:
        return ""
    html_doc = "\n<hr/>\n".join(text_chunks)
    # remove images entirely
    html_doc = re.sub(r"<img[^>]*>", "", html_doc, flags=re.IGNORECASE)
    text = md(html_doc).strip()
    return text if len(text) >= min_chars else ""


def render_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header = rows[0]
    body = rows[1:] or []
    def fmt_row(row):
        return "| " + " | ".join(cell or "" for cell in row) + " |"
    lines = [fmt_row(header), "| " + " | ".join("---" for _ in header) + " |"]
    for row in body:
        lines.append(fmt_row(row))
    return "\n".join(lines)


def ocr_pdf_with_pymupdf(path: str, lang: str, use_tesseract: bool = True, conf_threshold: float = 0.65):
    try:
        import fitz  # PyMuPDF
        from PIL import Image
        if not use_tesseract:
            return None
        import pytesseract
        import cv2
        import numpy as np
    except ImportError:
        return None

    doc = None
    elements: list = []
    try:
        doc = fitz.open(path)
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
            mode = "RGBA" if pix.alpha else "RGB"
            img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            if conf_threshold <= 0:
                text = pytesseract.image_to_string(img, lang=_join_lang(lang), config="--psm 6")
                if not text.strip():
                    text = pytesseract.image_to_string(img, lang=_join_lang(lang), config="--psm 4")
            else:
                arr = cv2.cvtColor(np.array(img), cv2.COLOR_RGBA2GRAY if mode == "RGBA" else cv2.COLOR_RGB2GRAY)
                arr = cv2.bilateralFilter(arr, 9, 75, 75)
                arr = cv2.adaptiveThreshold(arr, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 31, 2)
                text = ocr_text_from_array(arr, lang, conf_threshold)
            if text.strip():
                elements.append(SimpleNamespace(text=text))
    except Exception as err:
        print(f"[convert] PyMuPDF OCR failed for {path}: {err}", file=sys.stderr)
        return None
    finally:
        if doc is not None:
            doc.close()
    if not elements:
        elements.append(SimpleNamespace(text=""))
    return elements


def ocr_image(path: str, lang: str, use_tesseract: bool, conf_threshold: float) -> str:
    if not use_tesseract:
        return "[image omitted]"
    try:
        from PIL import Image
        import pytesseract
        import cv2
        import numpy as np
    except ImportError:
        return "[image omitted]"
    try:
        img = Image.open(path)
        if conf_threshold <= 0:
            text = pytesseract.image_to_string(img, lang=_join_lang(lang), config="--psm 6")
            if not text.strip():
                text = pytesseract.image_to_string(img, lang=_join_lang(lang), config="--psm 4")
        else:
            arr = cv2.cvtColor(np.array(img), cv2.COLOR_RGBA2GRAY if img.mode == "RGBA" else cv2.COLOR_RGB2GRAY)
            arr = cv2.bilateralFilter(arr, 9, 75, 75)
            arr = cv2.adaptiveThreshold(arr, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 31, 2)
            text = ocr_text_from_array(arr, lang, conf_threshold)
        return text.strip() or "[image omitted]"
    except Exception:
        return "[image omitted]"


def ocr_text_from_array(arr, lang: str, conf_threshold: float) -> str:
    try:
        import pytesseract
        from pytesseract import Output
    except ImportError:
        return ""
    if conf_threshold <= 0:
        return pytesseract.image_to_string(arr, lang=_join_lang(lang), config="--psm 6")
    try:
        threshold = max(0.0, min(conf_threshold, 0.99))
        data = pytesseract.image_to_data(
            arr,
            lang=_join_lang(lang),
            config="--psm 6",
            output_type=Output.DICT,
        )
        lines = defaultdict(list)
        order = []
        for i in range(len(data["text"])):
            conf = data["conf"][i]
            try:
                conf = float(conf)
            except Exception:
                continue
            if conf < (threshold * 100):
                continue
            txt = (data["text"][i] or "").strip()
            if not txt:
                continue
            key = (
                data.get("page_num", [0])[i],
                data.get("block_num", [0])[i],
                data.get("par_num", [0])[i],
                data.get("line_num", [0])[i],
            )
            if key not in lines:
                order.append(key)
            lines[key].append(txt)
        if not lines:
            return pytesseract.image_to_string(arr, lang=_join_lang(lang), config="--psm 6")
        kept_boxes = sum(len(v) for v in lines.values())
        total_boxes = sum(1 for t in data["text"] if t and t.strip())
        if total_boxes and kept_boxes / total_boxes < 0.3:
            return pytesseract.image_to_string(arr, lang=_join_lang(lang), config="--psm 6")
        sorted_keys = sorted(order)
        text_lines = [" ".join(lines[k]) for k in sorted_keys]
        return "\n".join(text_lines)
    except Exception:
        return pytesseract.image_to_string(arr, lang=_join_lang(lang), config="--psm 6")


def _parse_lang(lang: str | None) -> list[str]:
    if not lang:
        return ["eng"]
    if isinstance(lang, str):
        return [part.strip() for part in lang.replace("+", ",").split(",") if part.strip()]
    return lang


def _join_lang(lang: str | None) -> str:
    if not lang:
        return "eng"
    parts = _parse_lang(lang)
    return "+".join(parts)


def main():
    src = Path(ARGS.src).resolve()
    dst = Path(ARGS.dst).resolve()
    dst.mkdir(parents=True, exist_ok=True)

    files = list(iter_files(src, ARGS.glob))
    if not files:
        print("No documents found to convert.")
        return

    report_path = dst / 'convert_report.jsonl'
    ok = 0; skipped = 0; errors = 0
    with open(report_path, 'w', encoding='utf-8') as rep:
        for p in files:
            rec = convert_one(
                src,
                dst,
                p,
                overwrite=ARGS.overwrite,
                enable_ocr=ARGS.enable_ocr,
                ocr_lang=ARGS.ocr_lang,
                use_tesseract=ARGS.use_tesseract,
                ocr_conf_threshold=ARGS.ocr_conf_threshold,
                ocr_image_conf_threshold=ARGS.ocr_image_conf_threshold,
            )
            st = rec.get('status')
            if st == 'ok':
                ok += 1
            elif st == 'skipped':
                skipped += 1
            else:
                errors += 1
            rep.write(json.dumps(rec, ensure_ascii=False) + "\n")

    print(f"Converted: {ok}, skipped: {skipped}, errors: {errors}. Report: {report_path}")


def parse_args():
    ap = argparse.ArgumentParser()
    ap.add_argument('--src')
    ap.add_argument('--dst')
    ap.add_argument('--glob', nargs='*', default=[], help='Optional glob filters relative to src (e.g., **/*.pdf)')
    ap.add_argument('--overwrite', action='store_true')
    ap.add_argument('--enable-ocr', action='store_true', help='Force OCR pipeline for PDFs before partitioning')
    ap.add_argument('--ocr-lang', default='eng', help='OCR language codes (comma or + separated, e.g., eng+fra)')
    ap.add_argument('--use-tesseract', action='store_true', default=False, help='Enable Tesseract-based OCR (requires tesseract binary)')
    ap.add_argument('--ocr-conf-threshold', type=float, default=0.0, help='Confidence threshold (0-1) for OCR token filtering (0 disables filtering)')
    ap.add_argument('--ocr-image-conf-threshold', type=float, default=0.6, help='Confidence threshold for standalone images (0 disables filtering)')
    ap.add_argument('--from-config', action='store_true', help='Read convert_docs: tasks from RAG_CONFIG and run them')
    return ap.parse_args()


if __name__ == '__main__':
    ARGS = parse_args()
    if ARGS.from_config:
        cfg_path = os.environ.get('RAG_CONFIG', 'config.yaml')
        cfg = yaml.safe_load(open(cfg_path, 'r'))
        tasks = (cfg.get('convert_docs') or []) if isinstance(cfg, dict) else []
        if not tasks:
            print('[convert] No convert_docs tasks found in config.')
        for i, t in enumerate(tasks, 1):
            src = t.get('src'); dst = t.get('dst')
            gl = t.get('glob') or []
            ow = bool(t.get('overwrite', False))
            enable_ocr = bool(t.get('enable_ocr', False))
            ocr_lang = t.get('ocr_lang') or None
            use_tesseract = bool(t.get('use_tesseract', False))
            conf_threshold = float(t.get('ocr_conf_threshold', 0.0))
            image_conf_threshold = float(t.get('ocr_image_conf_threshold', 0.6))
            if not src or not dst:
                print(f"[convert] Task {i}: missing src/dst; skipping")
                continue
            src_p = Path(src).resolve(); dst_p = Path(dst).resolve()
            dst_p.mkdir(parents=True, exist_ok=True)
            files = list(iter_files(src_p, gl))
            print(f"[convert] Task {i}/{len(tasks)}: files={len(files)}  {src} -> {dst}")
            ok=0; skipped=0; err=0
            for p in files:
                rec = convert_one(
                    src_p,
                    dst_p,
                    p,
                    overwrite=ow,
                    enable_ocr=enable_ocr,
                    ocr_lang=ocr_lang,
                    use_tesseract=use_tesseract,
                    ocr_conf_threshold=conf_threshold,
                    ocr_image_conf_threshold=image_conf_threshold,
                )
                st = rec.get('status')
                if st == 'ok': ok += 1
                elif st == 'skipped': skipped += 1
                else: err += 1
            print(f"[convert] done: ok={ok} skipped={skipped} errors={err}")
    else:
        if not ARGS.src or not ARGS.dst:
            print('ERROR: --src and --dst required (or use --from-config)')
            raise SystemExit(2)
        main()
